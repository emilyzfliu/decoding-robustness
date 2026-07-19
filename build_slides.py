"""
Generate a Google-Slides-importable .pptx summarizing the decoding-robustness
findings (GPT-2 / OPT / Llama / Qwen). Import via: Slides > File > Import slides.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

HERE = "/home/varghese/decoding-robustness"
OUT = f"{HERE}/decoding_robustness_slides.pptx"
FIG_BEHAV = f"{HERE}/figures/compare/behavioral_compare.png"
FIG_REP = f"{HERE}/figures/compare/representation_compare.png"

# palette
INK = RGBColor(0x22, 0x22, 0x22)
MUTE = RGBColor(0x66, 0x66, 0x66)
ABS_C = RGBColor(0xD6, 0x60, 0x4D)   # absolute (red)
ROPE_C = RGBColor(0x21, 0x66, 0xAC)  # RoPE (blue)
BAND = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide():
    return prs.slides.add_slide(BLANK)

def textbox(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame

def set_run(r, text, size, bold=False, color=INK, italic=False):
    r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    r.font.name = "Arial"

def title(s, text, sub=None):
    tf = textbox(s, 0.6, 0.35, 12.1, 1.1)
    set_run(tf.paragraphs[0].add_run(), text, 30, True, INK)
    if sub:
        p = tf.add_paragraph(); set_run(p.add_run(), sub, 15, False, MUTE, True)
    # accent rule
    ln = s.shapes.add_shape(1, Inches(0.62), Inches(1.5), Inches(3.2), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = ROPE_C; ln.line.fill.background()

def bullets(s, items, l=0.7, t=1.75, w=6.0, h=5.2, size=16):
    tf = textbox(s, l, t, w, h)
    for i, (txt, lvl, *rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        color = rest[0] if rest else INK
        bold = rest[1] if len(rest) > 1 else False
        prefix = "" if lvl == 0 else "– "
        set_run(p.add_run(), ("• " if lvl == 0 else prefix) + txt, size - lvl, bold, color)
        p.space_after = Pt(6)

def image_fit(s, path, l, t, maxw, maxh):
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = maxw, maxw / ar
    if h > maxh:
        h, w = maxh, maxh * ar
    left = l + (maxw - w) / 2
    s.shapes.add_picture(path, Inches(left), Inches(t), Inches(w), Inches(h))

def caption(s, text, t):
    tf = textbox(s, 0.7, t, 12.0, 0.5)
    set_run(tf.paragraphs[0].add_run(), text, 12, False, MUTE, True)

def table(s, rows, l, t, w, h, header_fill=ROPE_C, col_w=None, fs=13, cell_style=None):
    # cell_style: {(row, col): (fill_rgb, text_rgb)} overrides for specific body cells
    cell_style = cell_style or {}
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    if col_w:
        for j, cw in enumerate(col_w):
            gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.margin_left = Inches(0.08); c.margin_top = Inches(0.03)
            c.margin_bottom = Inches(0.03)
            p = c.text_frame.paragraphs[0]
            txt_color = RGBColor(0xFF, 0xFF, 0xFF) if i == 0 else INK
            bold = i == 0
            if (i, j) in cell_style:
                fill_rgb, txt_color = cell_style[(i, j)]; bold = True
                c.fill.solid(); c.fill.fore_color.rgb = fill_rgb
            elif i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
            else:
                c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 else BAND
            set_run(p.add_run(), str(val), fs, bold, txt_color)
    return gt

# strippability cell colors
GRN_F, GRN_T = RGBColor(0xD9, 0xEA, 0xD3), RGBColor(0x2E, 0x7D, 0x32)
RED_F, RED_T = RGBColor(0xF6, 0xCF, 0xCB), RGBColor(0xA6, 0x1C, 0x00)

def panel(s, l, t, w, h, fill_rgb, border_rgb, head, lines):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill_rgb
    sh.line.color.rgb = border_rgb; sh.line.width = Pt(1.5)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22); tf.margin_top = Inches(0.16)
    set_run(tf.paragraphs[0].add_run(), head, 17, True, border_rgb)
    for txt, *rest in lines:
        p = tf.add_paragraph(); p.space_before = Pt(5)
        set_run(p.add_run(), txt, 14, rest[0] if rest else False, INK)
    return sh

# ---------------- Slide 1: title ----------------
s = slide()
band = s.shapes.add_shape(1, 0, Inches(2.5), SW, Inches(2.5))
band.fill.solid(); band.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFB); band.line.fill.background()
tf = textbox(s, 0.9, 2.7, 11.5, 1.4)
set_run(tf.paragraphs[0].add_run(),
        "Decoding Robustness: How Perturbations Propagate Through LMs", 34, True, INK)
p = tf.add_paragraph()
set_run(p.add_run(),
        "A representation-metric confound — and a cross-family (absolute vs RoPE) finding",
        18, False, ROPE_C)
tf2 = textbox(s, 0.9, 4.35, 11.5, 0.6)
set_run(tf2.paragraphs[0].add_run(),
        "GPT-2 (124M) · OPT-6.7B · Llama-3.1-8B · Qwen2.5-7B   |   WikiText-2, 100 passages × 128 tokens",
        14, False, MUTE)

# ---------------- Slide 2: setup ----------------
s = slide()
title(s, "Setup: three perturbations, four models",
      "Each clean passage run twice — clean vs perturbed — and compared")
bullets(s, [
    ("Perturbations, each isolating one interface:", 0, INK, True),
    ("Char substitution → tokenizer disruption", 1),
    ("Token substitution → lexical / semantic corruption", 1),
    ("Token shuffling → positional corruption", 1),
    ("Metrics:", 0, INK, True),
    ("Behavioral: perplexity, output divergence, logit-KL", 1),
    ("Representation: per-layer stripped cosine (content) + stripped CKA (geometry)", 1),
], w=6.4)
table(s, [
    ["Model", "Params", "Pos. encoding"],
    ["GPT-2", "124M", "learned absolute"],
    ["OPT-6.7B", "6.7B", "learned absolute"],
    ["Llama-3.1-8B", "8B", "RoPE"],
    ["Qwen2.5-7B", "7B", "RoPE"],
], l=7.5, t=2.0, w=5.2, h=2.8, col_w=[2.2, 1.2, 1.8])
caption(s, "Absolute vs RoPE at matched scale (~7–8B) is the controlled comparison.", 4.9)

# ---------------- Slide 3: behavioral ----------------
s = slide()
title(s, "Behavioral signatures differ by corruption type")
image_fit(s, FIG_BEHAV, 0.5, 1.5, 8.6, 5.5)
bullets(s, [
    ("Char: front-loaded — 5% already flips half the output", 0),
    ("Token: explosive perplexity (valid-but-wrong = max surprise)", 0),
    ("Shuffle: gentlest — models robust to local word-order", 0),
    ("Severity char < shuffle < token holds across all models", 0, ABS_C, True),
], l=9.2, t=1.9, w=3.9, size=14)

# ---------------- Slide 4: metric trap ----------------
s = slide()
title(s, "The representation metric was confounded",
      "Raw cosine 'recovered' to ~0.95 at the final layer even when outputs were totally different")
bullets(s, [
    ("Cause: massive activations (few 'rogue' dims implementing attention sinks)", 0),
    ("Near-constant, huge magnitude → survive perturbation → fake 'similarity'", 1),
    ("Every off-the-shelf metric is fooled — differently:", 0, INK, True),
], w=12.0, t=1.7)
table(s, [
    ["Metric", "Weights dims by", "Fooled by", "Result"],
    ["Raw cosine", "magnitude (+offset)", "the huge offset", "falsely HIGH (0.95)"],
    ["Z-scored cosine", "1 / std", "flat dims' noise", "falsely LOW (0.04)"],
    ["Plain CKA", "variance", "preserved positional variance", "falsely HIGH (0.999)"],
], l=0.7, t=3.15, w=11.9, h=1.9, col_w=[2.3, 2.9, 3.6, 3.1])
bullets(s, [
    ("Fix: drop the top-k highest-variance dims per layer, THEN measure", 0, ABS_C, True),
], t=5.3, w=12.0)

# ---------------- Slide 5: cross-family representation ----------------
s = slide()
title(s, "RoPE keeps a similarity that absolute models lose")
image_fit(s, FIG_REP, 0.4, 1.5, 9.0, 5.6)
bullets(s, [
    ("Top (CKA): RoPE (dashed) high plateau; absolute (solid) low", 0, ROPE_C, True),
    ("Bottom (cosine): inverts — absolute higher", 0, ABS_C, True),
    ("RoPE's high CKA = pure positional geometry — the confound itself", 0),
], l=9.5, t=2.0, w=3.6, size=13)

# ---------------- Slide 6: the finding ----------------
s = slide()
title(s, "Two separable effects: positional encoding × depth",
      "After stripping the outlier dims, do the two metrics agree?   (cosine = content, CKA = geometry)")
G, R = (GRN_F, GRN_T), (RED_F, RED_T)
conv_txt, gap_txt = "cosine ≈ CKA", "cosine ≪ CKA"
table(s, [
    ["Model", "Pos. encoding", "Early layers", "Late layers"],
    ["GPT-2 (124M)", "absolute", conv_txt, conv_txt],
    ["OPT-6.7B", "absolute", conv_txt, gap_txt],
    ["Llama-3.1-8B", "RoPE", gap_txt, gap_txt],
    ["Qwen2.5-7B", "RoPE", gap_txt, gap_txt],
], l=0.7, t=2.05, w=9.6, h=2.5, col_w=[2.6, 2.2, 2.4, 2.4], fs=14,
   cell_style={
       (1, 2): G, (1, 3): G,
       (2, 2): G, (2, 3): R,
       (3, 2): R, (3, 3): R,
       (4, 2): R, (4, 3): R,
   })
caption(s, "cosine ≈ CKA → metrics converge (confound gone)      cosine ≪ CKA → gap persists (confound remains)", 4.6)
# "how to read it" callouts
bullets(s, [
    ("Read DOWN ‘Early layers’ ↓ — absolute (OPT) converge; RoPE (Llama, Qwen) don’t", 0, ROPE_C, True),
    ("→ positional encoding controls whether the metrics converge early", 1, INK, True),
    ("Read ACROSS each row → — the gap re-opens with depth, even for absolute OPT", 0, ABS_C, True),
    ("→ a depth effect, independent of positional encoding", 1, INK, True),
    ("So a persistent gap is NOT unique to RoPE — OPT (no RoPE) shows it late", 0, ABS_C, True),
], t=5.0, w=12.6, size=14)

# ---------------- Slide 7: mechanism ----------------
s = slide()
title(s, "Why does this happen?")
# premise box (full width)
prem = s.shapes.add_shape(1, Inches(0.7), Inches(1.7), Inches(11.9), Inches(1.5))
prem.fill.solid(); prem.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
prem.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); prem.line.width = Pt(1); prem.shadow.inherit = False
ptf = prem.text_frame; ptf.word_wrap = True
ptf.margin_left = Inches(0.25); ptf.margin_top = Inches(0.14)
set_run(ptf.paragraphs[0].add_run(),
        "Char corruption changes the tokens — but keeps their POSITIONS fixed.", 17, True, INK)
p = ptf.add_paragraph(); p.space_before = Pt(6)
set_run(p.add_run(),
        "So the positional signal is identical in clean vs perturbed. A metric reading geometry (CKA) "
        "latches onto that shared signal → looks ‘similar’; a metric reading direction/content "
        "(cosine) sees the changed tokens → looks ‘different’.", 14, False, INK)
# two contrast panels
panel(s, 0.7, 3.5, 5.8, 2.7, RGBColor(0xFB, 0xEC, 0xE9), ABS_C,
      "Absolute position  (GPT-2, OPT)",
      [("Position is stored in a FEW dimensions.",),
       ("Drop those dims → the shared positional signal is gone",),
       ("→ CKA falls to meet cosine.",),
       ("Result: cosine and CKA CONVERGE (early layers)", True)])
panel(s, 6.8, 3.5, 5.8, 2.7, RGBColor(0xE8, 0xF0, 0xF7), ROPE_C,
      "RoPE  (Llama, Qwen)",
      [("Position is spread across ALL dimensions (rotations).",),
       ("There is no small set of dims to drop",),
       ("→ CKA stays pinned high (~0.8).",),
       ("Result: the cosine–CKA GAP PERSISTS", True)])
caption(s, "With depth, position spreads out even in absolute models → the gap re-opens in their late layers too.", 6.4)

# ---------------- Slide 8: takeaways ----------------
s = slide()
title(s, "Takeaways & open questions")
bullets(s, [
    ("Takeaways", 0, ROPE_C, True),
    ("Off-the-shelf representation-similarity metrics are confounded by massive", 1),
    ("activations — always strip outlier dims before interpreting", 1),
    ("The strip budget k is model-specific (k=5 for GPT-2; wide models need more,", 1),
    ("or never converge under RoPE)", 1),
    ("Absolute vs RoPE produce qualitatively different geometry under perturbation", 1),
    ("Open questions", 0, ABS_C, True),
    ("Disentangle depth vs positional encoding at matched depth", 1),
    ("Phase 2: attention / head-ablation (GQA complicates single-head study)", 1),
    ("Char-perplexity non-monotonicity; k-sensitivity sweep for publication", 1),
], w=12.4, t=1.75, size=16)

prs.save(OUT)
print("wrote", OUT, f"({len(prs.slides._sldIdLst)} slides)")
